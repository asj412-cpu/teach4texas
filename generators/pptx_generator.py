"""Generate PowerPoint/Google Slides versions of task card products.

Requires: pip install python-pptx
"""
import json
import io
import os

from .brand import (
    OUTPUT_DIR, CONTENT_DIR, NAVY, BURNT, GREEN, GOLD, WHITE,
    ensure_output_dir
)
from . import visual_aids

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


SUBJECT_COLORS_RGB = {
    "math": RGBColor(191, 87, 0) if HAS_PPTX else None,
    "rla": RGBColor(84, 130, 53) if HAS_PPTX else None,
    "science": RGBColor(191, 135, 0) if HAS_PPTX else None,
}

NAVY_RGB = RGBColor(27, 54, 93) if HAS_PPTX else None
WHITE_RGB = RGBColor(255, 255, 255) if HAS_PPTX else None


def _generate_visual_bytes(spec):
    """Generate a visual aid as PNG bytes."""
    vtype = spec["type"]
    params = spec.get("params", {})

    generators = {
        "number_line": visual_aids.draw_number_line,
        "fraction_bar": visual_aids.draw_fraction_bar,
        "area_model": visual_aids.draw_area_model,
        "coordinate_grid": visual_aids.draw_coordinate_grid,
        "bar_graph": visual_aids.draw_bar_graph,
        "data_table": visual_aids.draw_data_table,
        "geometric_shape": visual_aids.draw_geometric_shape,
        "graphic_organizer": visual_aids.draw_graphic_organizer,
        "food_web": visual_aids.draw_food_web,
        "force_arrows": visual_aids.draw_force_arrows,
    }

    gen = generators.get(vtype)
    if gen is None:
        return None

    img = gen(**params)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _add_title_slide(prs, data):
    """Add a title slide to the presentation."""
    accent = SUBJECT_COLORS_RGB.get(data["subject"], NAVY_RGB)
    subject_labels = {"math": "Math", "rla": "RLA", "science": "Science"}
    subj = subject_labels.get(data["subject"], data["subject"].upper())

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_RGB

    # Title
    left, top = Inches(0.8), Inches(1.5)
    width, height = Inches(8.4), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"STAAR {subj} Task Cards"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE_RGB
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    left, top = Inches(1.5), Inches(3.5)
    width, height = Inches(7), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Grade {data['grade']} | 32 TEKS-Aligned Questions"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE_RGB
    p.alignment = PP_ALIGN.CENTER

    # Brand
    left, top = Inches(3), Inches(5.5)
    width, height = Inches(4), Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Teach4"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE_RGB
    run2 = p.add_run()
    run2.text = "Texas"
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER


def _add_question_slide(prs, q, subject):
    """Add a single question slide."""
    accent = SUBJECT_COLORS_RGB.get(subject, NAVY_RGB)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Card number + TEKS header bar
    left, top = Inches(0.3), Inches(0.3)
    width, height = Inches(9.4), Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"  #{q['id']}  "
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = NAVY_RGB
    run2 = p.add_run()
    run2.text = f"  TEKS {q['teks']}"
    run2.font.size = Pt(14)
    run2.font.color.rgb = accent

    # Question text
    has_visual = q.get("visual") is not None
    q_top = Inches(1.2)
    q_height = Inches(1.8) if has_visual else Inches(3)

    left = Inches(0.5)
    width = Inches(9)
    txBox = slide.shapes.add_textbox(left, q_top, width, q_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = q["question"]
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(45, 45, 45)

    # Visual aid
    if has_visual:
        img_buf = _generate_visual_bytes(q["visual"])
        if img_buf:
            slide.shapes.add_picture(
                img_buf, Inches(2.5), Inches(3.2), height=Inches(2))

    # Answer choices
    choices_top = Inches(5.5) if has_visual else Inches(4.5)
    for i, choice in enumerate(q["choices"]):
        row, col = divmod(i, 2)
        left = Inches(1) + col * Inches(4.5)
        top = choices_top + row * Inches(0.6)
        txBox = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        letter = choice[0]
        text = choice[2:].strip()
        run = p.add_run()
        run.text = f"{letter})  "
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = NAVY_RGB
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(45, 45, 45)


def generate_task_cards_pptx(json_path, output_name=None):
    """Generate a PPTX file from task card content JSON.

    Args:
        json_path: Path to the content JSON
        output_name: Optional filename without extension
    """
    if not HAS_PPTX:
        print("  SKIPPED (python-pptx not installed)")
        return None

    with open(json_path) as f:
        data = json.load(f)

    subject = data["subject"]
    grade = data["grade"]

    if output_name is None:
        output_name = f"task-cards-{subject}-grade{grade}"

    print(f"Generating PPTX: {output_name}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, data)

    for q in data["questions"]:
        _add_question_slide(prs, q, subject)

    output_path = os.path.join(OUTPUT_DIR, f"{output_name}.pptx")
    ensure_output_dir()
    prs.save(output_path)
    size = os.path.getsize(output_path)
    print(f"  PPTX: {output_path} ({size:,} bytes)")
    return output_path


def generate_all_pptx():
    """Generate PPTX versions for all task card content files."""
    if not HAS_PPTX:
        print("python-pptx not installed. Run: pip install python-pptx")
        return []

    ensure_output_dir()
    generated = []

    for fname in sorted(os.listdir(CONTENT_DIR)):
        if fname.endswith("_tasks.json"):
            json_path = os.path.join(CONTENT_DIR, fname)
            path = generate_task_cards_pptx(json_path)
            if path:
                generated.append(path)

    print(f"\nGenerated {len(generated)} PPTX file(s)")
    return generated


if __name__ == "__main__":
    generate_all_pptx()
